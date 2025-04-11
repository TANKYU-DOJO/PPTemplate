import os
import re
import tempfile
import subprocess
import openpyxl
import pptx
from pptx.dml.color import RGBColor

def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    p = paragraph._p
    for idx, run in enumerate(paragraph.runs):
        if idx > 0:
            p.remove(run._r)

    paragraph.runs[0].text = new_text

def replace_text(slide, before, after):
    for shp in slide.shapes:
        if shp.has_text_frame:
            if before in shp.text:
                new_text = re.sub(before, after, shp.text)
                replace_paragraph_text_retaining_initial_formatting(shp.text_frame.paragraphs[0], new_text)

def generate(path_template: str, path_data: str, path_output: str, sheetname: str | None = None) -> None:
    wb = openpyxl.load_workbook(path_data, data_only=True)
    ws = wb.active
    if sheetname is not None:
        ws = wb[sheetname]

    keys = []
    for cell in ws[1]:
        if cell.value is None: break
        keys.append(str(cell.value))

    with tempfile.TemporaryDirectory() as td:
        counter = 1
        while ws[counter + 1][0].value is not None:
            template = pptx.Presentation(path_template)
            for i, key in enumerate(keys):
                value = ws[counter + 1][i].value
                if value is None:
                    for shape in template.slides[0].shapes:
                        if hasattr(shape, 'text') and shape.text == key: # 図形の文字列がキーと一致していた場合
                            shape.fill.background() # 見えなくする
                    replace_text(template.slides[0], key, '')
                else:
                    value = str(value)
                    if re.fullmatch(r'#[0-9a-zA-Z]{6}', value) is not None: # カラーコードで指定されていた場合
                        for shape in template.slides[0].shapes:
                            if hasattr(shape, 'text') and shape.text == key:
                                hex = value.lstrip("#")
                                rgb = [int(hex[i:i+2], 16) for i in range(0, 6, 2)]
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                                shape.text = ''
                    elif re.fullmatch(r'\./.*\.[0-9a-zA-Z]+', value) is not None: # ./(任意の文字列).(拡張子) の場合、画像ファイルとして処理する
                        for shape in template.slides[0].shapes:
                            if hasattr(shape, 'text') and shape.text == key:
                                path = os.path.join(os.path.dirname(path_data), value.lstrip('./'))
                                template.slides[0].shapes.add_picture(path, shape.left, shape.top, shape.width, shape.height)
                                # shapeを削除
                                XML_reference = shape._sp
                                XML_reference.getparent().remove(XML_reference)
                    else: replace_text(template.slides[0], key, value) # 通常の文字列として処理
            template.save(os.path.join(td, str(counter) + '.pptx'))
            counter += 1

        command = [os.path.join(os.path.dirname(__file__), 'mergepptx'), '-o', path_output]
        for i in range(1,counter):
            command.append(os.path.join(td, str(i) + '.pptx'))
        subprocess.run(command)

    wb.close()